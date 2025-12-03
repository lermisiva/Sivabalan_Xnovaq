from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Function to add a slide with title and content
def add_slide(prs, title, content_items):
    slide_layout = prs.slide_layouts[1] # Bullet layout
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    
    # Set Title
    title_shape = shapes.title
    title_shape.text = title
    
    # Set Content
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, item in enumerate(content_items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(18)

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # Title layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]
title.text = "ShopSphere Data Engineering & Analytics Project"
subtitle.text = "Based on BRD v1.0\nBuilding a Scalable, Unified Data Platform"

# Slide 2: Client Overview
add_slide(prs, "Client Overview: ShopSphere Inc.", [
    "Domain: US-based e-commerce retailer operating in 6 regions (Northeast, Midwest, etc.).",
    "Scale: Serves 150K+ daily visitors and processes 12K orders daily.",
    "Current Tech: On-premise PostgreSQL, Microservices logs, and Vendor CSVs.",
    "Status: Rapid growth hindered by scattered data and manual reporting."
])

# Slide 3: Pain Point 1 - Fragmented Data
add_slide(prs, "The Problem: Fragmented Data Landscape", [
    "Data is siloed across incompatible systems:",
    "- PostgreSQL: Transactional data (hard to scale).",
    "- S3 Buckets: Unstructured event logs.",
    "- Vendor CSVs: Inventory updates (manual & delayed).",
    "Impact: Teams spend hours manually pulling and merging data."
])

# Slide 4: Pain Point 2 - Operational Bottlenecks
add_slide(prs, "Operational Bottlenecks & Manual Work", [
    "No Real-Time Insights: Business relies on yesterday's reports.",
    "Reactive: When products go viral (e.g., TikTok), inventory runs out before teams notice.",
    "Inefficiency: Reporting team spends 70% of their time collecting data instead of analyzing it.",
    "Blind Spots: No unified 'Customer 360' view to track buying patterns."
])

# Slide 5: Business Impact
add_slide(prs, "Business Impact", [
    "Revenue Loss: Inventory discrepancies lead to Out-of-Stock items and refunds.",
    "Customer Trust: Confusion over availability damages the brand.",
    "Scalability Failures: During Black Friday/Christmas, systems slow down and pipelines fail.",
    "Delayed Decisions: Lack of real-time data prevents agile marketing."
])

# Slide 6: Proposed Solution
add_slide(prs, "Proposed Solution: Unified Data Platform", [
    "Objective: Build a scalable, automated, cloud-based data platform.",
    "Single Source of Truth: Unify transactional, log, and vendor data.",
    "Automation: Replace manual Excel work with automated pipelines.",
    "Scalability: Handle 4x traffic spikes (Holiday/Black Friday) without failure."
])

# Slide 7: Key Capabilities
add_slide(prs, "Key Capabilities & Deliverables", [
    "Real-Time Visibility: Instant insights into inventory and trends.",
    "Customer 360: Unified view of orders, returns, and clickstream behavior.",
    "Automated Dashboards: KPI monitoring to replace manual reports.",
    "High Reliability: Low latency systems ensuring data quality."
])

# Save the presentation
prs.save('ShopSphere_Project_Proposal.pptx')
print("Presentation saved as 'ShopSphere_Project_Proposal.pptx'")